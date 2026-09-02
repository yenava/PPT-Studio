#!/usr/bin/env python3
from __future__ import annotations

import argparse
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from textppt_common import page_dirs, read_json, write_json


def validate_deck(run_dir: str | Path, pptx_path: str | Path) -> dict:
    run_dir = Path(run_dir).expanduser().resolve()
    pptx_path = Path(pptx_path).expanduser().resolve()
    errors: list[str] = []
    page_reports: list[dict] = []
    if not pptx_path.exists():
        report = {"schema_version": 1, "passed": False, "errors": [f"Missing PPTX: {pptx_path}"], "pages": []}
        write_json(run_dir / "validation.json", report)
        return report
    if not zipfile.is_zipfile(pptx_path):
        errors.append("Output is not a valid ZIP-based PPTX package.")
    try:
        presentation = Presentation(pptx_path)
    except Exception as exc:
        report = {"schema_version": 1, "passed": False, "errors": errors + [f"PPTX open failed: {exc}"], "pages": []}
        write_json(run_dir / "validation.json", report)
        return report

    directories = page_dirs(run_dir)
    if len(presentation.slides) != len(directories):
        errors.append(f"Slide count {len(presentation.slides)} does not match page count {len(directories)}.")

    for index, page_dir in enumerate(directories):
        page_errors: list[str] = []
        ocr = read_json(page_dir / "ocr.json")
        expected_entries = ocr.get("entries", [])
        check = read_json(page_dir / "clean-base-check.json")
        if not check.get("passed"):
            page_errors.append("Clean-base OCR check did not pass.")
        with Image.open(page_dir / "source.png") as source, Image.open(page_dir / "clean_base.png") as clean:
            if source.size != clean.size:
                page_errors.append(f"Clean-base dimensions {clean.size} do not match source {source.size}.")
        if index >= len(presentation.slides):
            page_errors.append("Slide is missing from the PPTX.")
            page_reports.append({"page_id": page_dir.name, "passed": False, "errors": page_errors})
            continue
        slide = presentation.slides[index]
        pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        text_boxes = [shape for shape in slide.shapes if shape.has_text_frame]
        if len(pictures) != 1:
            page_errors.append(f"Expected exactly one slide picture; found {len(pictures)}.")
        elif any(
            abs(value) > 2
            for value in (
                pictures[0].left,
                pictures[0].top,
                pictures[0].width - presentation.slide_width,
                pictures[0].height - presentation.slide_height,
            )
        ):
            page_errors.append("The clean-base picture does not fill the slide.")
        if len(text_boxes) != len(expected_entries):
            page_errors.append(f"Expected {len(expected_entries)} text boxes; found {len(text_boxes)}.")
        actual_text = [shape.text for shape in text_boxes]
        expected_text = [str(entry.get("text", "")) for entry in expected_entries]
        if actual_text != expected_text:
            page_errors.append("Text-box contents or order do not match OCR entries.")
        for shape in text_boxes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > presentation.slide_width + 2 or shape.top + shape.height > presentation.slide_height + 2:
                page_errors.append(f"Text box is outside slide bounds: {shape.text[:40]!r}.")
        low_confidence = [
            entry["id"] for entry in expected_entries if float(entry.get("confidence", 1.0)) < 0.5
        ]
        page_reports.append(
            {
                "page_id": page_dir.name,
                "passed": not page_errors,
                "errors": page_errors,
                "picture_count": len(pictures),
                "text_box_count": len(text_boxes),
                "low_confidence_entries": low_confidence,
                "clean_base_attempt": check.get("attempt"),
            }
        )
        errors.extend(f"{page_dir.name}: {message}" for message in page_errors)

    report = {
        "schema_version": 1,
        "passed": not errors,
        "pptx": str(pptx_path),
        "slide_count": len(presentation.slides),
        "errors": errors,
        "pages": page_reports,
    }
    write_json(run_dir / "validation.json", report)
    return report


def main() -> int:
    parser = argparse.ArgumentParser(description="Validate an editable-text PPTX against its run artifacts.")
    parser.add_argument("run_dir")
    parser.add_argument("pptx")
    args = parser.parse_args()
    report = validate_deck(args.run_dir, args.pptx)
    print(f"Validation {'passed' if report['passed'] else 'failed'} with {len(report['errors'])} error(s).")
    return 0 if report["passed"] else 2


if __name__ == "__main__":
    raise SystemExit(main())
