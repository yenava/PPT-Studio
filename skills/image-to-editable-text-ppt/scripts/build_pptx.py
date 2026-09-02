#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from textppt_common import page_dirs, read_json


ALIGNMENT = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


def _text_units(value: str) -> float:
    return sum(1.0 if ord(character) > 255 else 0.55 for character in value)


def _fitted_font_size(entry: dict, box_width_pt: float, box_height_pt: float) -> float:
    requested = float(entry.get("font_pt", 18))
    lines = str(entry.get("text", "")).splitlines() or [""]
    widest = max((_text_units(line) for line in lines), default=1.0)
    width_scale = box_width_pt / max(1.0, widest * requested)
    height_scale = box_height_pt / max(1.0, len(lines) * requested * 1.15)
    return round(max(6.0, min(requested, requested * width_scale * 0.96, requested * height_scale * 0.96)), 2)


def build_deck(run_dir: str | Path, output: str | Path, *, font_name: str = "PingFang SC") -> Path:
    run_dir = Path(run_dir).expanduser().resolve()
    manifest = read_json(run_dir / "run.json")
    slide_width_in = float(manifest["slide_width_in"])
    slide_height_in = float(manifest["slide_height_in"])
    presentation = Presentation()
    presentation.slide_width = Inches(slide_width_in)
    presentation.slide_height = Inches(slide_height_in)
    blank = presentation.slide_layouts[6]

    for page_dir in page_dirs(run_dir):
        source = page_dir / "source.png"
        clean_base = page_dir / "clean_base.png"
        if not clean_base.exists():
            raise FileNotFoundError(f"Missing clean base: {clean_base}")
        check = read_json(page_dir / "clean-base-check.json")
        if not check.get("passed"):
            raise ValueError(f"Clean-base check has not passed: {page_dir}")
        ocr = read_json(page_dir / "ocr.json")
        source_width = float(ocr["source"]["width_px"])
        source_height = float(ocr["source"]["height_px"])
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(clean_base),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )
        for entry in ocr.get("entries", []):
            x, y, width, height = (float(value) for value in entry["box_px"])
            left = Inches(x / source_width * slide_width_in)
            top = Inches(y / source_height * slide_height_in)
            box_width = Inches(width / source_width * slide_width_in)
            box_height = Inches(height / source_height * slide_height_in)
            shape = slide.shapes.add_textbox(left, top, box_width, box_height)
            frame = shape.text_frame
            frame.clear()
            frame.word_wrap = True
            frame.auto_size = MSO_AUTO_SIZE.NONE
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            frame.margin_left = 0
            frame.margin_right = 0
            frame.margin_top = 0
            frame.margin_bottom = 0
            paragraph = frame.paragraphs[0]
            paragraph.alignment = ALIGNMENT.get(str(entry.get("alignment", "left")), PP_ALIGN.LEFT)
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            paragraph.line_spacing = 1.0
            run = paragraph.add_run()
            run.text = str(entry.get("text", ""))
            run.font.name = font_name
            font_size = _fitted_font_size(entry, box_width.pt, box_height.pt)
            run.font.size = Pt(font_size)
            color = str(entry.get("color_hex", "000000")).lstrip("#")
            if len(color) != 6:
                color = "000000"
            run.font.color.rgb = RGBColor.from_string(color.upper())

    output = Path(output).expanduser().resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    return output


def main() -> int:
    parser = argparse.ArgumentParser(description="Build a PPTX from clean bases and OCR text layers.")
    parser.add_argument("run_dir")
    parser.add_argument("--out", default="output.pptx")
    parser.add_argument("--font", default="PingFang SC")
    args = parser.parse_args()
    output = Path(args.out)
    if not output.is_absolute():
        output = Path(args.run_dir) / output
    result = build_deck(args.run_dir, output, font_name=args.font)
    print(f"Wrote {result}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
