#!/usr/bin/env python3
"""Build an image-based 16:9 PPTX from confirmed slide images."""

from __future__ import annotations

import argparse
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}


def collect_images(paths: list[Path]) -> list[Path]:
    images: list[Path] = []
    for path in paths:
        if path.is_dir():
            images.extend(p for p in path.iterdir() if p.suffix.lower() in IMAGE_EXTS)
        elif path.suffix.lower() in IMAGE_EXTS:
            images.append(path)
    return sorted(images, key=lambda p: p.name)


def image_box(image_path: Path, slide_w: int, slide_h: int, fit: str) -> tuple[int, int, int, int]:
    if fit == "stretch":
        return 0, 0, slide_w, slide_h

    with Image.open(image_path) as img:
        img_w, img_h = img.size
    img_ratio = img_w / img_h
    slide_ratio = slide_w / slide_h

    if fit == "contain":
        if img_ratio >= slide_ratio:
            w = slide_w
            h = int(slide_w / img_ratio)
        else:
            h = slide_h
            w = int(slide_h * img_ratio)
    else:  # cover
        if img_ratio >= slide_ratio:
            h = slide_h
            w = int(slide_h * img_ratio)
        else:
            w = slide_w
            h = int(slide_w / img_ratio)

    x = int((slide_w - w) / 2)
    y = int((slide_h - h) / 2)
    return x, y, w, h


def build(images: list[Path], out: Path, fit: str) -> None:
    if not images:
        raise SystemExit("No slide images found.")

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for image in images:
        slide = prs.slides.add_slide(blank)
        x, y, w, h = image_box(image, prs.slide_width, prs.slide_height, fit)
        slide.shapes.add_picture(str(image), x, y, width=w, height=h)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)


def main() -> int:
    parser = argparse.ArgumentParser(description="Create a 16:9 image-based PPTX from slide images.")
    parser.add_argument("--images", nargs="+", required=True, type=Path, help="Image files or directories.")
    parser.add_argument("--out", required=True, type=Path)
    parser.add_argument("--fit", choices=["cover", "contain", "stretch"], default="cover")
    args = parser.parse_args()

    build(collect_images(args.images), args.out, args.fit)
    print(args.out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
